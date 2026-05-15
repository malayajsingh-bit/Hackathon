import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from utils.config import BRAND, FONT_FAMILY, FOOTER_TEXT, SLIDE_DIMENSIONS


_DEFAULT_TEMPLATE = os.path.join(os.path.dirname(os.path.dirname(__file__)), "templates", "indiamart_default.pptx")

# Slide geometry constants (inches) — content slides only
# Title slide uses its own hardcoded values (1.57 / 3.09) — see add_title_slide()
_TITLE_TOP    = 0.70
_TITLE_HEIGHT = 0.80
_CONTENT_TOP  = 1.86
_CONTENT_H    = 4.24   # 6.10 - 1.86
_CALLOUT_TOP  = 6.10
_CALLOUT_H    = 0.60
_FOOTER_TOP   = 6.85
_LEFT_MARGIN  = 0.70
_SLIDE_W      = 11.80  # usable width (13.33 - 2×0.70)


class PPTRenderer:
    def __init__(self, profile: dict):
        if os.path.exists(_DEFAULT_TEMPLATE):
            self.prs = Presentation(_DEFAULT_TEMPLATE)
            self._clear_slides()
        else:
            self.prs = Presentation()
            self.prs.slide_width = SLIDE_DIMENSIONS["width"]
            self.prs.slide_height = SLIDE_DIMENSIONS["height"]
        self.profile = profile
        self.vis = profile["visual_preferences"]
        self.min_font = self.vis["font_size_minimum"]
        self.max_bullets = self.vis.get("max_bullet_points_per_slide", 8)

    # ------------------------------------------------------------------ #
    #  Template setup helpers
    # ------------------------------------------------------------------ #

    def _clear_slides(self):
        """Remove all slides, preserving the slide master/theme/colors."""
        from pptx.oxml.ns import qn
        sldIdLst = self.prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn('r:id'))
            sldIdLst.remove(sldId)
            try:
                self.prs.part.drop_rel(rId)
            except Exception:
                pass

    def _blank_layout(self):
        """Find blank layout by name; fall back to index 6."""
        for layout in self.prs.slide_layouts:
            if layout.name.lower() in ('blank', 'blank slide'):
                return layout
        try:
            return self.prs.slide_layouts[6]
        except IndexError:
            return self.prs.slide_layouts[-1]

    def _title_slide_layout(self):
        """Return the Title Slide layout (index 0) — Indiamart-designed first page."""
        for layout in self.prs.slide_layouts:
            if layout.name.lower() in ('title slide', 'title_slide', 'title'):
                return layout
        return self.prs.slide_layouts[0]

    def _strip_placeholders(self, slide):
        """Remove every placeholder element from slide XML."""
        for ph in list(slide.placeholders):
            ph._element.getparent().remove(ph._element)

    # ------------------------------------------------------------------ #
    #  Public API
    # ------------------------------------------------------------------ #

    def add_title_slide(self, title: str, subtitle: str, date: str):
        slide = self.prs.slides.add_slide(self._title_slide_layout())

        # First-slide geometry (hardcoded — independent of content-slide constants)
        # Title  : Y=1.57"  (aligned with grey line in template), font 32pt
        # Content: Y=3.09"  (below grey line), one textbox
        # X width: textbox center at (slide_center - 3) keeps right half free for logo
        #   slide_center = 13.333/2 = 6.667"  →  center = 3.667"
        #   left = 0.70"  →  width = 2×(3.667 - 0.70) = 5.93"
        _TS_Y_TITLE   = 1.57
        _TS_Y_CONTENT = 3.09
        _TS_FONT      = 32
        _TS_LEFT      = _LEFT_MARGIN                        # 0.70"
        _TS_WIDTH     = 2 * (13.333 / 2 - 3.0 - _LEFT_MARGIN)  # ~5.93"

        # Strip ALL placeholders — we place everything ourselves so positions
        # are exact and no template placeholder overrides our layout.
        for ph in list(slide.placeholders):
            ph._element.getparent().remove(ph._element)

        # ONE title textbox — word_wrap on, wraps to 2nd line if long, never splits
        self._add_textbox(slide, title,
                          left=Inches(_TS_LEFT), top=Inches(_TS_Y_TITLE),
                          width=Inches(_TS_WIDTH), height=Inches(1.6),
                          font_size=Pt(_TS_FONT), color=BRAND["dark"],
                          bold=True, alignment=PP_ALIGN.LEFT)

        # ONE subtitle/date textbox below the grey line
        sub_text = f"{subtitle}   {date}" if subtitle else date
        self._add_textbox(slide, sub_text,
                          left=Inches(_TS_LEFT), top=Inches(_TS_Y_CONTENT),
                          width=Inches(_TS_WIDTH), height=Inches(0.8),
                          font_size=Pt(18), color=BRAND["dark"],
                          alignment=PP_ALIGN.LEFT)

    def add_content_slide(self, slide_content: dict, chart_path: str = None,
                          diagram_path: str = None):
        slide_type = slide_content.get("slide_type", "content")

        if slide_type == "executive_summary":
            self._add_summary_slide(slide_content)
        elif slide_type == "chart" or chart_path:
            self._add_chart_slide(slide_content, chart_path)
        elif slide_type == "diagram" or diagram_path:
            self._add_diagram_slide(slide_content, diagram_path)
        elif slide_type == "comparison":
            self._add_comparison_slide(slide_content)
        elif slide_type == "ask":
            self._add_ask_slide(slide_content)
        else:
            self._add_bullet_slide(slide_content)

    # ------------------------------------------------------------------ #
    #  Slide type renderers
    # ------------------------------------------------------------------ #

    def _add_summary_slide(self, content: dict):
        slide = self._create_base_slide(content.get("title", "Executive Summary"))

        bullets = content.get("bullets", [])[:self.max_bullets]
        metrics = [b for b in bullets if any(c in b for c in ["₹", "%", "+", "-", "×", "x"])]
        non_metrics = [b for b in bullets if b not in metrics]

        # TEXT CONTENT first (top), then rectangular metric boxes below
        bullet_h = (_CONTENT_H * 0.52) if metrics else _CONTENT_H
        if non_metrics:
            self._add_bullets(slide, non_metrics,
                              left=Inches(_LEFT_MARGIN), top=Inches(_CONTENT_TOP),
                              width=Inches(_SLIDE_W), height=Inches(bullet_h))

        if metrics:
            n = len(metrics)
            card_top = _CONTENT_TOP + (bullet_h + 0.15 if non_metrics else 0.0)
            box_w = min(3.4, _SLIDE_W / n)
            if n == 1:
                self._add_metric_card(slide, metrics[0],
                                      Inches(_LEFT_MARGIN), Inches(card_top),
                                      Inches(box_w), Inches(1.6))
            else:
                gap = (_SLIDE_W - n * box_w) / (n - 1)
                for i, metric in enumerate(metrics):
                    left = Inches(_LEFT_MARGIN + i * (box_w + gap))
                    self._add_metric_card(slide, metric, left,
                                          Inches(card_top), Inches(box_w), Inches(1.6))

        self._add_callout(slide, content.get("key_callout", ""))
        self._add_speaker_notes(slide, content.get("speaker_notes", ""))

    def _add_bullet_slide(self, content: dict):
        slide = self._create_base_slide(content.get("title", ""))

        bullets = content.get("bullets", [])[:self.max_bullets]
        if bullets:
            self._add_bullets(slide, bullets,
                              left=Inches(_LEFT_MARGIN), top=Inches(_CONTENT_TOP),
                              width=Inches(_SLIDE_W), height=Inches(_CONTENT_H))

        self._add_callout(slide, content.get("key_callout", ""))
        self._add_speaker_notes(slide, content.get("speaker_notes", ""))

    def _add_chart_slide(self, content: dict, chart_path: str):
        slide = self._create_base_slide(content.get("title", ""))

        if chart_path and os.path.exists(chart_path):
            slide.shapes.add_picture(
                chart_path,
                left=Inches(0.5), top=Inches(_CONTENT_TOP),
                width=Inches(7.8), height=Inches(_CONTENT_H))

            bullets = content.get("bullets", [])[:self.max_bullets]
            if bullets:
                self._add_bullets(slide, bullets,
                                  left=Inches(8.6), top=Inches(_CONTENT_TOP + 0.2),
                                  width=Inches(4.4), height=Inches(_CONTENT_H - 0.4))
        else:
            bullets = content.get("bullets", [])[:self.max_bullets]
            if bullets:
                self._add_bullets(slide, bullets,
                                  left=Inches(_LEFT_MARGIN), top=Inches(_CONTENT_TOP),
                                  width=Inches(_SLIDE_W), height=Inches(_CONTENT_H))

        self._add_callout(slide, content.get("key_callout", ""))
        self._add_speaker_notes(slide, content.get("speaker_notes", ""))

    def _add_diagram_slide(self, content: dict, diagram_path: str):
        slide = self._create_base_slide(content.get("title", ""))

        if diagram_path and os.path.exists(diagram_path):
            # Use generous margins so diagram breathes — not edge-to-edge
            slide.shapes.add_picture(
                diagram_path,
                left=Inches(1.2), top=Inches(_CONTENT_TOP + 0.1),
                width=Inches(11.0), height=Inches(_CONTENT_H - 0.2))
        else:
            self._add_textbox(slide, "[Diagram placeholder]",
                              left=Inches(2), top=Inches(3.5), width=Inches(9), height=Inches(1),
                              font_size=Pt(18), color=BRAND["text_muted"],
                              alignment=PP_ALIGN.CENTER)

        self._add_speaker_notes(slide, content.get("speaker_notes", ""))

    def _add_comparison_slide(self, content: dict):
        slide = self._create_base_slide(content.get("title", ""))

        bullets = content.get("bullets", [])[:self.max_bullets]
        mid = len(bullets) // 2
        left_items = bullets[:mid] if mid > 0 else bullets
        right_items = bullets[mid:] if mid > 0 else []

        col_w = 5.9
        label_h = 0.45
        col_top = _CONTENT_TOP + 0.05
        col_h = _CONTENT_H - 0.1
        label_font = Pt(max(self.min_font - 6, 13))

        # Only draw a box when the column actually has content
        if left_items:
            shape_left = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(col_top), Inches(col_w), Inches(col_h))
            shape_left.fill.solid()
            shape_left.fill.fore_color.rgb = RGBColor(0xFF, 0xE4, 0xE4)
            shape_left.line.color.rgb = BRAND["danger"]
            self._add_textbox(slide, "BEFORE / CURRENT",
                              left=Inches(0.8), top=Inches(col_top + 0.1),
                              width=Inches(col_w - 0.5), height=Inches(label_h),
                              font_size=label_font, color=BRAND["danger"], bold=True)
            self._add_bullets(slide, left_items,
                              left=Inches(0.8), top=Inches(col_top + label_h + 0.25),
                              width=Inches(col_w - 0.5),
                              height=Inches(col_h - label_h - 0.4))

        if right_items:
            shape_right = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.9), Inches(col_top), Inches(col_w), Inches(col_h))
            shape_right.fill.solid()
            shape_right.fill.fore_color.rgb = RGBColor(0xE4, 0xFF, 0xE4)
            shape_right.line.color.rgb = BRAND["success"]
            self._add_textbox(slide, "AFTER / PROPOSED",
                              left=Inches(7.2), top=Inches(col_top + 0.1),
                              width=Inches(col_w - 0.5), height=Inches(label_h),
                              font_size=label_font, color=BRAND["success"], bold=True)
            self._add_bullets(slide, right_items,
                              left=Inches(7.2), top=Inches(col_top + label_h + 0.25),
                              width=Inches(col_w - 0.5),
                              height=Inches(col_h - label_h - 0.4))

        self._add_callout(slide, content.get("key_callout", ""))
        self._add_speaker_notes(slide, content.get("speaker_notes", ""))

    def _add_ask_slide(self, content: dict):
        slide = self._create_base_slide(content.get("title", "Decision Needed"))

        bullets = content.get("bullets", [])[:self.max_bullets]
        if bullets:
            # All decision points in ONE textbox — no background box
            self._add_bullets(slide, bullets,
                              left=Inches(_LEFT_MARGIN), top=Inches(_CONTENT_TOP),
                              width=Inches(_SLIDE_W), height=Inches(_CONTENT_H),
                              color=BRAND["dark"])

        self._add_speaker_notes(slide, content.get("speaker_notes", ""))

    # ------------------------------------------------------------------ #
    #  Base slide & shared primitives
    # ------------------------------------------------------------------ #

    def _create_base_slide(self, title: str):
        slide = self.prs.slides.add_slide(self._blank_layout())
        self._strip_placeholders(slide)

        self._add_textbox(slide, title,
                          left=Inches(_LEFT_MARGIN), top=Inches(_TITLE_TOP),
                          width=Inches(_SLIDE_W), height=Inches(_TITLE_HEIGHT),
                          font_size=Pt(max(self.min_font, 26)),
                          color=BRAND["dark"], bold=True, alignment=PP_ALIGN.LEFT)

        self._add_footer(slide)
        return slide

    def _add_bullets(self, slide, bullets: list, left, top, width, height, color=None):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Dynamic font size: shrink slightly if there are many bullets
        n = len(bullets)
        base_size = max(self.min_font, 18)
        if n > 5:
            base_size = max(self.min_font - 2, 15)
        if n > 7:
            base_size = max(self.min_font - 4, 13)

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(4)
            p.space_after = Pt(4)
            p.text = ""

            run_bullet = p.add_run()
            run_bullet.text = "•  "
            run_bullet.font.color.rgb = BRAND["primary"]
            run_bullet.font.size = Pt(base_size)
            run_bullet.font.name = FONT_FAMILY

            run_text = p.add_run()
            run_text.text = bullet
            run_text.font.color.rgb = color or BRAND["text"]
            run_text.font.size = Pt(base_size)
            run_text.font.name = FONT_FAMILY

    def _add_metric_card(self, slide, metric_text: str, left, top, width, height):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
        shape.line.color.rgb = BRAND["primary"]
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(10)
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)

        p = tf.paragraphs[0]
        p.text = metric_text
        p.font.size = Pt(max(self.min_font - 4, 15))
        p.font.color.rgb = BRAND["dark"]
        p.font.name = FONT_FAMILY
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def _add_callout(self, slide, text: str):
        if not text:
            return

        txBox = slide.shapes.add_textbox(
            Inches(_LEFT_MARGIN), Inches(_CALLOUT_TOP),
            Inches(_SLIDE_W), Inches(_CALLOUT_H))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        run_label = p.add_run()
        run_label.text = "KEY INSIGHT:  "
        run_label.font.bold = True
        run_label.font.size = Pt(max(self.min_font - 6, 13))
        run_label.font.color.rgb = BRAND["dark"]
        run_label.font.name = FONT_FAMILY

        run_text = p.add_run()
        run_text.text = text
        run_text.font.size = Pt(max(self.min_font - 6, 13))
        run_text.font.color.rgb = BRAND["text"]
        run_text.font.name = FONT_FAMILY

        p.alignment = PP_ALIGN.LEFT

    def _add_brand_bar(self, slide):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Pt(6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BRAND["accent"]
        bar.line.fill.background()

    def _add_footer(self, slide):
        self._add_textbox(slide, FOOTER_TEXT,
                          left=Inches(_LEFT_MARGIN), top=Inches(_FOOTER_TOP),
                          width=Inches(9), height=Inches(0.38),
                          font_size=Pt(9), color=BRAND["text_muted"],
                          alignment=PP_ALIGN.LEFT)

        slide_num = len(self.prs.slides)
        self._add_textbox(slide, str(slide_num),
                          left=Inches(12.1), top=Inches(_FOOTER_TOP),
                          width=Inches(0.8), height=Inches(0.38),
                          font_size=Pt(9), color=BRAND["text_muted"],
                          alignment=PP_ALIGN.RIGHT)

    def _add_textbox(self, slide, text: str, left, top, width, height,
                     font_size=Pt(14), color=None, bold=False,
                     alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color or BRAND["text"]
        p.font.bold = bold
        p.font.name = FONT_FAMILY
        p.alignment = alignment

    def _add_speaker_notes(self, slide, notes: str):
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def save(self, output_path: str):
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
        self.prs.save(output_path)
        return output_path
