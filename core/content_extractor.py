import os
import json
import subprocess
import tempfile
import shutil
import requests
import pandas as pd
from utils.claude_client import ClaudeClient

SYSTEM_PROMPT = """You are a business analyst at Indiamart Intermesh Ltd.
Your job is to extract structured content from source material for presentations.
Extract ONLY what exists in the source. Do NOT invent or hallucinate data.
If data is ambiguous, flag it with [VERIFY] prefix."""

EXTRACTION_PROMPT = """Analyze the following source material and extract structured content for a presentation.

SOURCE MATERIAL:
{content}

ADDITIONAL CONTEXT FROM USER:
{context}

Extract into this exact JSON structure:
{{
  "title_suggestion": "short presentation title (max 10 words)",
  "key_metrics": [
    {{"metric": "name", "value": "number/amount", "change": "+/-X%", "period": "time period"}}
  ],
  "findings": ["finding 1", "finding 2"],
  "achievements": ["achievement 1", "achievement 2"],
  "challenges": ["challenge 1", "challenge 2"],
  "recommendations": ["recommendation 1", "recommendation 2"],
  "decisions_needed": ["decision 1"],
  "chart_data": [
    {{
      "title": "chart title",
      "type": "bar",
      "labels": ["label1", "label2"],
      "datasets": [{{"label": "series name", "values": [100, 200]}}]
    }}
  ],
  "diagram_suggestions": [
    {{
      "type": "architecture|flowchart|comparison|timeline",
      "title": "diagram title",
      "description": "what the diagram should show",
      "nodes": ["node1", "node2"],
      "connections": ["node1 -> node2"]
    }}
  ]
}}

Rules:
- Extract actual numbers from the source, do not make up metrics
- chart_data values must be numeric (no strings like "₹100")
- If no chart data is available, return empty chart_data array
- diagram_suggestions should only be included when content naturally suits a visual
- Keep findings/achievements/challenges concise (1 sentence each)"""


def read_pdf(file_path: str) -> str:
    import fitz
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    doc.close()
    return text


def read_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def read_csv(file_path: str) -> str:
    df = pd.read_csv(file_path)
    summary = f"CSV Data: {len(df)} rows, {len(df.columns)} columns\n"
    summary += f"Columns: {', '.join(df.columns.tolist())}\n\n"
    summary += "First 20 rows:\n"
    summary += df.head(20).to_string(index=False)
    numeric_cols = df.select_dtypes(include="number").columns
    if len(numeric_cols) > 0:
        summary += "\n\nNumeric Summary:\n"
        summary += df[numeric_cols].describe().to_string()
    return summary


def read_excel(file_path: str) -> str:
    df = pd.read_excel(file_path)
    summary = f"Excel Data: {len(df)} rows, {len(df.columns)} columns\n"
    summary += f"Columns: {', '.join(df.columns.tolist())}\n\n"
    summary += "First 20 rows:\n"
    summary += df.head(20).to_string(index=False)
    numeric_cols = df.select_dtypes(include="number").columns
    if len(numeric_cols) > 0:
        summary += "\n\nNumeric Summary:\n"
        summary += df[numeric_cols].describe().to_string()
    return summary


def read_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def read_pptx(file_path: str) -> str:
    """Extract text content from an existing PowerPoint file."""
    from pptx import Presentation
    prs = Presentation(file_path)
    content_parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(row_data))
        if slide_texts:
            content_parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                content_parts.append(f"  Speaker Notes: {notes}")

    summary = f"PowerPoint: {len(prs.slides)} slides\n\n"
    summary += "\n\n".join(content_parts)
    return summary


def read_image(file_path: str) -> str:
    """Return a placeholder description for image files. Claude API with vision can enhance this."""
    import base64
    ext = os.path.splitext(file_path)[1].lower()
    mime_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp"}
    mime = mime_map.get(ext, "image/png")
    size = os.path.getsize(file_path)
    size_kb = size / 1024
    return (f"--- Image: {os.path.basename(file_path)} ---\n"
            f"Type: {mime}, Size: {size_kb:.1f} KB\n"
            f"[Image uploaded — will be embedded directly into the presentation]\n"
            f"File path: {file_path}")


def read_url(url: str) -> str:
    """Fetch and extract text content from a webpage."""
    try:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        resp = requests.get(url.strip(), headers=headers, timeout=15)
        resp.raise_for_status()
        html = resp.text

        # Simple HTML to text extraction (no BeautifulSoup dependency)
        import re
        # Remove script and style tags
        html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
        html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL | re.IGNORECASE)
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', ' ', html)
        # Clean whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        # Decode HTML entities
        text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
        text = text.replace('&nbsp;', ' ').replace('&quot;', '"')

        if len(text) > 20000:
            text = text[:20000] + "\n\n[TRUNCATED]"

        return f"--- Webpage: {url} ---\n{text}"
    except Exception as e:
        return f"Error fetching URL {url}: {e}"


def read_google_sheet_csv(url: str) -> str:
    """Read a Google Sheet by converting the share URL to CSV export format."""
    try:
        # Convert Google Sheets URL to CSV export URL
        import re
        match = re.search(r'/spreadsheets/d/([a-zA-Z0-9_-]+)', url)
        if not match:
            return "Error: Could not parse Google Sheets URL. Use the share link format."

        sheet_id = match.group(1)
        # Extract gid if present
        gid_match = re.search(r'gid=(\d+)', url)
        gid = gid_match.group(1) if gid_match else "0"

        csv_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid={gid}"

        resp = requests.get(csv_url, timeout=15)
        resp.raise_for_status()

        from io import StringIO
        df = pd.read_csv(StringIO(resp.text))
        summary = f"Google Sheet Data: {len(df)} rows, {len(df.columns)} columns\n"
        summary += f"Columns: {', '.join(df.columns.tolist())}\n\n"
        summary += "First 30 rows:\n"
        summary += df.head(30).to_string(index=False)
        numeric_cols = df.select_dtypes(include="number").columns
        if len(numeric_cols) > 0:
            summary += "\n\nNumeric Summary:\n"
            summary += df[numeric_cols].describe().to_string()
        return summary
    except Exception as e:
        return f"Error reading Google Sheet: {e}\nMake sure the sheet is shared with 'Anyone with the link'."


FILE_READERS = {
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".csv": read_csv,
    ".xlsx": read_excel,
    ".xls": read_excel,
    ".txt": read_text,
    ".md": read_text,
    ".pptx": read_pptx,
    ".ppt": read_pptx,
    ".png": read_image,
    ".jpg": read_image,
    ".jpeg": read_image,
    ".gif": read_image,
    ".webp": read_image,
    ".bmp": read_image,
}


def read_github_repo(repo_url: str) -> str:
    """Clone a GitHub repo and read key files (README, docs, source structure)."""
    tmp_dir = tempfile.mkdtemp(prefix="ghrepo_")
    try:
        clean_url = repo_url.strip().rstrip("/")
        if not clean_url.endswith(".git"):
            clean_url += ".git"

        subprocess.run(
            ["git", "clone", "--depth", "1", clean_url, tmp_dir],
            capture_output=True, text=True, timeout=60,
            cwd=tempfile.gettempdir()
        )

        content_parts = []

        # Read README
        for readme_name in ["README.md", "readme.md", "README.txt", "README"]:
            readme_path = os.path.join(tmp_dir, readme_name)
            if os.path.exists(readme_path):
                with open(readme_path, "r", encoding="utf-8", errors="ignore") as f:
                    content_parts.append(f"--- README ---\n{f.read()[:10000]}")
                break

        # Read directory structure
        tree_lines = []
        for root, dirs, files in os.walk(tmp_dir):
            dirs[:] = [d for d in dirs if d not in [".git", "node_modules", "venv", "__pycache__", ".next", "dist", "build"]]
            level = root.replace(tmp_dir, "").count(os.sep)
            if level > 3:
                continue
            indent = "  " * level
            folder_name = os.path.basename(root)
            tree_lines.append(f"{indent}{folder_name}/")
            sub_indent = "  " * (level + 1)
            for f in files[:15]:
                tree_lines.append(f"{sub_indent}{f}")

        content_parts.append(f"--- Directory Structure ---\n" + "\n".join(tree_lines[:200]))

        # Read key source files (limited)
        key_patterns = ["package.json", "requirements.txt", "pyproject.toml",
                        "Dockerfile", "docker-compose.yml", ".env.example",
                        "CHANGELOG.md", "CONTRIBUTING.md"]
        for pattern in key_patterns:
            fpath = os.path.join(tmp_dir, pattern)
            if os.path.exists(fpath):
                with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()[:3000]
                    content_parts.append(f"--- {pattern} ---\n{text}")

        # Read up to 5 source files for context
        source_extensions = [".py", ".js", ".ts", ".java", ".go", ".rb"]
        source_count = 0
        for root, dirs, files in os.walk(tmp_dir):
            dirs[:] = [d for d in dirs if d not in [".git", "node_modules", "venv", "__pycache__"]]
            for f in files:
                if source_count >= 5:
                    break
                ext = os.path.splitext(f)[1]
                if ext in source_extensions:
                    fpath = os.path.join(root, f)
                    rel_path = os.path.relpath(fpath, tmp_dir)
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as fh:
                        text = fh.read()[:3000]
                        content_parts.append(f"--- Source: {rel_path} ---\n{text}")
                        source_count += 1

        # Read git log for recent activity
        try:
            log_result = subprocess.run(
                ["git", "log", "--oneline", "-20"],
                capture_output=True, text=True, timeout=10, cwd=tmp_dir
            )
            if log_result.stdout:
                content_parts.append(f"--- Recent Commits ---\n{log_result.stdout}")
        except Exception:
            pass

        return "\n\n".join(content_parts)
    except Exception as e:
        return f"Error reading GitHub repo: {e}"
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def read_openproject_tickets(base_url: str, api_key: str, project_id: str = None) -> str:
    """Fetch work packages (tickets) from OpenProject API."""
    headers = {
        "Authorization": f"Basic {api_key}",
        "Content-Type": "application/json",
    }

    content_parts = []

    # Fetch work packages
    endpoint = f"{base_url.rstrip('/')}/api/v3/work_packages"
    params = {"pageSize": 50, "sortBy": '[["updatedAt","desc"]]'}
    if project_id:
        endpoint = f"{base_url.rstrip('/')}/api/v3/projects/{project_id}/work_packages"

    try:
        resp = requests.get(endpoint, headers=headers, params=params, timeout=30)
        resp.raise_for_status()
        data = resp.json()

        total = data.get("total", 0)
        content_parts.append(f"--- OpenProject Tickets (Total: {total}) ---\n")

        tickets_by_status = {}
        tickets_by_type = {}
        tickets_by_priority = {}

        for wp in data.get("_embedded", {}).get("elements", []):
            subject = wp.get("subject", "")
            status = wp.get("_links", {}).get("status", {}).get("title", "Unknown")
            wp_type = wp.get("_links", {}).get("type", {}).get("title", "Unknown")
            priority = wp.get("_links", {}).get("priority", {}).get("title", "Unknown")
            assignee = wp.get("_links", {}).get("assignee", {}).get("title", "Unassigned")
            updated = wp.get("updatedAt", "")[:10]
            description = wp.get("description", {}).get("raw", "")[:200] if wp.get("description") else ""

            tickets_by_status[status] = tickets_by_status.get(status, 0) + 1
            tickets_by_type[wp_type] = tickets_by_type.get(wp_type, 0) + 1
            tickets_by_priority[priority] = tickets_by_priority.get(priority, 0) + 1

            content_parts.append(
                f"[{status}] [{wp_type}] [{priority}] {subject}\n"
                f"  Assignee: {assignee} | Updated: {updated}\n"
                f"  {description}\n"
            )

        # Summary stats
        summary = "\n--- Ticket Summary ---\n"
        summary += "By Status:\n" + "\n".join([f"  {k}: {v}" for k, v in tickets_by_status.items()])
        summary += "\nBy Type:\n" + "\n".join([f"  {k}: {v}" for k, v in tickets_by_type.items()])
        summary += "\nBy Priority:\n" + "\n".join([f"  {k}: {v}" for k, v in tickets_by_priority.items()])
        content_parts.insert(1, summary)

    except requests.exceptions.RequestException as e:
        content_parts.append(f"Error fetching OpenProject data: {e}")
        content_parts.append("Make sure the URL, API key, and project ID are correct.")

    return "\n".join(content_parts)


class ContentExtractor:
    def __init__(self, claude_client: ClaudeClient):
        self.claude = claude_client

    def extract_from_github(self, repo_url: str, context: str = "") -> dict:
        repo_content = read_github_repo(repo_url)
        return self._extract(repo_content, context)

    def extract_from_openproject(self, base_url: str, api_key: str,
                                  project_id: str = None, context: str = "") -> dict:
        ticket_content = read_openproject_tickets(base_url, api_key, project_id)
        return self._extract(ticket_content, context)

    def extract_from_url(self, url: str, context: str = "") -> dict:
        page_content = read_url(url)
        return self._extract(page_content, context)

    def extract_from_google_sheet(self, url: str, context: str = "") -> dict:
        sheet_content = read_google_sheet_csv(url)
        return self._extract(sheet_content, context)

    def extract_from_previous_ppt(self, file_path: str, new_data_paths: list = None,
                                   new_data_text: str = "", context: str = "") -> dict:
        """Extract from an existing PPT and optionally merge with new data."""
        old_content = read_pptx(file_path)
        parts = [f"--- EXISTING PRESENTATION (to update) ---\n{old_content}"]

        if new_data_paths:
            for path in new_data_paths:
                ext = os.path.splitext(path)[1].lower()
                reader = FILE_READERS.get(ext)
                if reader:
                    try:
                        text = reader(path)
                        parts.append(f"--- NEW DATA: {os.path.basename(path)} ---\n{text}")
                    except Exception as e:
                        parts.append(f"--- NEW DATA: {os.path.basename(path)} (ERROR: {e}) ---")

        if new_data_text:
            parts.append(f"--- NEW DATA (user provided) ---\n{new_data_text}")

        update_context = ("UPDATE THIS PRESENTATION with the new data provided. "
                          "Keep the same structure and slide flow, but refresh all metrics, "
                          "findings, and charts with the new data. " + (context or ""))

        combined = "\n\n".join(parts)
        return self._extract(combined, update_context)

    def get_image_paths(self, file_paths: list) -> list:
        """Filter and return paths that are images (for direct embedding in PPT)."""
        image_exts = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
        return [p for p in file_paths if os.path.splitext(p)[1].lower() in image_exts]

    def extract_from_files(self, file_paths: list, context: str = "") -> dict:
        all_content = []
        for path in file_paths:
            ext = os.path.splitext(path)[1].lower()
            reader = FILE_READERS.get(ext)
            if reader:
                try:
                    text = reader(path)
                    all_content.append(f"--- File: {os.path.basename(path)} ---\n{text}")
                except Exception as e:
                    all_content.append(f"--- File: {os.path.basename(path)} (ERROR: {e}) ---")
        combined = "\n\n".join(all_content)
        return self._extract(combined, context)

    def extract_from_text(self, text: str, context: str = "") -> dict:
        return self._extract(text, context)

    def extract_from_topic(self, topic: str, context: str = "") -> dict:
        prompt = f"Topic: {topic}\n\nAdditional context: {context}" if context else f"Topic: {topic}"
        return self._extract(prompt, context)

    def _extract(self, content: str, context: str) -> dict:
        if len(content) > 50000:
            content = content[:50000] + "\n\n[TRUNCATED - content too long]"

        prompt = EXTRACTION_PROMPT.format(content=content, context=context)
        result = self.claude.generate_json(SYSTEM_PROMPT, prompt)

        for key in ["key_metrics", "findings", "achievements", "challenges",
                     "recommendations", "decisions_needed", "chart_data", "diagram_suggestions"]:
            if key not in result:
                result[key] = []
        if "title_suggestion" not in result:
            result["title_suggestion"] = "Presentation"

        return result
