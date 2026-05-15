"""
Converts a .pptx file to a list of PNG slide preview images.
Strategy:
  1. Try LibreOffice headless (best fidelity, uses the real renderer)
  2. Fall back to a Pillow-based renderer that reads python-pptx shapes
"""
import os
import subprocess
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont

PREVIEW_W = 1280   # output pixel width
PREVIEW_H = 720    # output pixel height (16:9)


# ─── Public API ──────────────────────────────────────────────────────────────

def generate_previews(pptx_path: str, output_dir: str) -> list:
    """Return a list of PNG paths, one per slide."""
    os.makedirs(output_dir, exist_ok=True)
    paths = _try_libreoffice(pptx_path, output_dir)
    if paths:
        return paths
    return _render_with_pillow(pptx_path, output_dir)


# ─── LibreOffice path ─────────────────────────────────────────────────────────

def _try_libreoffice(pptx_path: str, output_dir: str) -> list:
    candidates = [
        "libreoffice", "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    for cmd in candidates:
        try:
            r = subprocess.run(
                [cmd, "--headless", "--convert-to", "png",
                 "--outdir", output_dir, pptx_path],
                capture_output=True, timeout=120, text=True
            )
            if r.returncode == 0:
                found = sorted([
                    os.path.join(output_dir, f)
                    for f in os.listdir(output_dir)
                    if f.lower().startswith(base.lower()) and f.lower().endswith(".png")
                ])
                if found:
                    return found
        except Exception:
            continue
    return []


# ─── Pillow fallback renderer ─────────────────────────────────────────────────

def _load_font(size_pt: float, bold: bool = False, scale: float = 1.0) -> ImageFont.FreeTypeFont:
    size_px = max(9, int(size_pt * scale))
    win_fonts = r"C:\Windows\Fonts"
    candidates = (
        [r"calibrib.ttf", r"arialbd.ttf"] if bold
        else [r"calibri.ttf", r"arial.ttf", r"DejaVuSans.ttf"]
    )
    for name in candidates:
        path = os.path.join(win_fonts, name)
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size_px)
            except Exception:
                pass
    return ImageFont.load_default()


def _rgb(pptx_color) -> tuple:
    """Extract (r, g, b) from a pptx RGBColor or return grey."""
    try:
        v = int(pptx_color)
        return ((v >> 16) & 0xFF, (v >> 8) & 0xFF, v & 0xFF)
    except Exception:
        return (31, 41, 55)


def _render_with_pillow(pptx_path: str, output_dir: str) -> list:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)

    # Scale: map slide EMUs to PREVIEW pixels
    sc_x = PREVIEW_W / prs.slide_width.emu
    sc_y = PREVIEW_H / prs.slide_height.emu
    sc   = min(sc_x, sc_y)
    W    = int(prs.slide_width.emu  * sc)
    H    = int(prs.slide_height.emu * sc)
    # pt-to-px scale for fonts: 1pt = 1/72 inch; rendered at ~96 px/inch
    font_sc = sc * prs.slide_width.emu / PREVIEW_W * 1.18

    paths = []
    for idx, slide in enumerate(prs.slides):
        img  = Image.new("RGB", (W, H), "#FFFFFF")
        draw = ImageDraw.Draw(img)

        # Draw a light brand top-bar (Indiamart orange accent)
        draw.rectangle([0, 0, W, max(4, int(6 * sc))], fill="#FF6B35")

        for shape in slide.shapes:
            try:
                left = int(shape.left   * sc)
                top  = int(shape.top    * sc)
                w    = int(shape.width  * sc)
                h    = int(shape.height * sc)

                # ── Picture shapes ──────────────────────────────────────────
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        thumb = Image.open(BytesIO(shape.image.blob)).convert("RGB")
                        thumb = thumb.resize((max(1, w), max(1, h)), Image.LANCZOS)
                        img.paste(thumb, (left, top))
                    except Exception:
                        pass
                    continue

                # ── Filled background shapes ────────────────────────────────
                if hasattr(shape, "fill"):
                    try:
                        if shape.fill.type is not None:
                            fc = shape.fill.fore_color.rgb
                            r, g, b = _rgb(fc)
                            opacity = 200
                            overlay = Image.new("RGBA", (max(1, w), max(1, h)),
                                                (r, g, b, opacity))
                            img.paste(overlay, (left, top),
                                      mask=overlay.split()[3])
                    except Exception:
                        pass

                # ── Text frames ─────────────────────────────────────────────
                if not shape.has_text_frame:
                    continue

                y_cur = top + max(5, int(6 * sc))
                for para in shape.text_frame.paragraphs:
                    raw = para.text.strip()
                    if not raw:
                        y_cur += max(4, int(6 * sc))
                        continue

                    # Font attributes from first run
                    pt, bold, color = 14, False, (31, 41, 55)
                    for run in para.runs:
                        if run.font.size:
                            pt = run.font.size.pt
                        bold = bool(run.font.bold)
                        try:
                            color = _rgb(run.font.color.rgb)
                        except Exception:
                            pass
                        break

                    font   = _load_font(pt, bold, font_sc)
                    line_h = max(10, int(pt * font_sc * 1.35))
                    max_x  = left + w - max(8, int(8 * sc))

                    # Word-wrap
                    words = raw.split()
                    line  = ""
                    for word in words:
                        test = (line + " " + word).strip()
                        try:
                            bb = draw.textbbox((left, y_cur), test, font=font)
                            tw = bb[2]
                        except Exception:
                            tw = left + len(test) * int(pt * font_sc * 0.6)
                        if tw <= max_x or not line:
                            line = test
                        else:
                            if y_cur < top + h - line_h:
                                draw.text((left + max(4, int(6 * sc)), y_cur),
                                          line, fill=color, font=font)
                            y_cur += line_h
                            line = word
                    if line and y_cur < top + h - line_h:
                        draw.text((left + max(4, int(6 * sc)), y_cur),
                                  line, fill=color, font=font)
                    y_cur += line_h

            except Exception:
                continue

        # Slide number badge (bottom-right)
        badge_font = _load_font(11, False, font_sc * 0.8)
        draw.text((W - int(32 * sc), H - int(20 * sc)),
                  str(idx + 1), fill="#9CA3AF", font=badge_font)

        # Light border
        draw.rectangle([0, 0, W - 1, H - 1], outline="#E5E7EB", width=2)

        out = os.path.join(output_dir, f"preview_{idx + 1:03d}.png")
        img.save(out, "PNG", optimize=True)
        paths.append(out)

    return paths
